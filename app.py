from flask import Flask, render_template, request, send_file
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.dml.color import RGBColor

app = Flask(__name__)

# Load menu items from Excel
excel_file = "menu_items.xlsx"  # Ensure this file exists
pptx_template = "menu_template.pptx"  # Your existing PowerPoint template
menu_df = pd.read_excel(excel_file, sheet_name="data")
menu_df = menu_df.astype(str)  # Convert all columns to strings

# Convert menu items for display
menu_items = menu_df[['index', 'items']].to_dict(orient='records')

def replace_text_in_pptx(selected_data):
    """Replaces placeholders in an existing PowerPoint template with selected meal data."""
    prs = Presentation(pptx_template)
    today_date = datetime.today().strftime("%A, %B %d, %Y")
    day_name = datetime.today().strftime("%A")
    
    print("\n--- Debug: Replacing Text in PowerPoint Placeholders ---")

    for slide_number, slide in enumerate(prs.slides, start=1):
        print(f"\n🔍 Checking Slide {slide_number}...")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.runs:
                        continue  # Skip empty paragraphs

                    full_text = "".join(run.text for run in paragraph.runs)  # Get full paragraph text

                    # Replace date and day placeholders
                    full_text = full_text.replace("{{date}}", today_date)
                    full_text = full_text.replace("{{day}}", day_name)

                    for meal, details in selected_data.items():
                        for i, (item, cal, allergen) in enumerate(zip(details['items'], details['calories'], details['allergens']), start=1):
                            item_placeholder = f"{{{{{meal}_{i:02}}}}}"  # Example: {{breakfast_01}}
                            cal_placeholder = f"{{{{{meal}_c_{i:02}}}}}"  # Example: {{breakfast_c_01}}
                            allergen_placeholder = f"{{{{{meal}_a_{i:02}}}}}"  # Example: {{breakfast_a_01}}

                            if item_placeholder in full_text:
                                print(f"Replacing {item_placeholder} with {item}")
                                full_text = full_text.replace(item_placeholder, item)

                            if cal_placeholder in full_text:
                                print(f"Replacing {cal_placeholder} with {cal}")
                                full_text = full_text.replace(cal_placeholder, cal)

                            if allergen_placeholder in full_text:
                                print(f"Replacing {allergen_placeholder} with {allergen}")
                                full_text = full_text.replace(allergen_placeholder, allergen)

                                # Change allergen text color if not "Free / no allergic ingredients."
                                if allergen != "Free / no allergic ingredients.":
                                    for run in paragraph.runs:
                                        if allergen_placeholder in run.text:
                                            run.font.color.rgb = RGBColor(255, 0, 0)  # Red color

                    for run in paragraph.runs:
                        run.text = ""  # Clear existing text
                    paragraph.runs[0].text = full_text  # Set updated text in first run

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            if not paragraph.runs:
                                continue  # Skip empty paragraphs

                            full_text = "".join(run.text for run in paragraph.runs)
                            full_text = full_text.replace("{{date}}", today_date)
                            full_text = full_text.replace("{{day}}", day_name)

                            for meal, details in selected_data.items():
                                for i, (item, cal, allergen) in enumerate(zip(details['items'], details['calories'], details['allergens']), start=1):
                                    item_placeholder = f"{{{{{meal}_{i:02}}}}}"
                                    cal_placeholder = f"{{{{{meal}_c_{i:02}}}}}"
                                    allergen_placeholder = f"{{{{{meal}_a_{i:02}}}}}"

                                    if item_placeholder in full_text:
                                        full_text = full_text.replace(item_placeholder, item)

                                    if cal_placeholder in full_text:
                                        full_text = full_text.replace(cal_placeholder, cal)

                                    if allergen_placeholder in full_text:
                                        full_text = full_text.replace(allergen_placeholder, allergen)

                                        # Change allergen text color if not "Free / no allergic ingredients."
                                        for run in paragraph.runs:
                                            if allergen_placeholder in run.text and allergen != "Free / no allergic ingredients.":
                                                run.font.color.rgb = RGBColor(255, 0, 0)  # Red color

                            for run in paragraph.runs:
                                run.text = ""  # Clear existing text
                            paragraph.runs[0].text = full_text

    output_pptx = f"{today_date}_menu.pptx"
    prs.save(output_pptx)
    print("--- Debug: PowerPoint Updated Successfully ---\n")
    return output_pptx

@app.route('/', methods=['GET', 'POST'])
def index():
    selected_data = {}

    if request.method == 'POST':
        selected_indices = {
            'breakfast': [request.form.get(f'breakfast_{i}') for i in range(1, 8)],
            'lunch': [request.form.get(f'lunch_{i}') for i in range(1, 9)],
            'dinner': [request.form.get(f'dinner_{i}') for i in range(1, 9)]
        }

        for meal, indexes in selected_indices.items():
            selected_data[meal] = {
                'items': [menu_df.loc[menu_df['index'] == idx, 'items'].values[0] if idx and idx in menu_df['index'].values else "" for idx in indexes],
                'calories': [menu_df.loc[menu_df['index'] == idx, 'calories'].values[0] if idx and idx in menu_df['index'].values else "" for idx in indexes],
                'allergens': [menu_df.loc[menu_df['index'] == idx, 'allergens'].values[0] if idx and idx in menu_df['index'].values else "" for idx in indexes],
            }

        updated_pptx = replace_text_in_pptx(selected_data)
        return send_file(updated_pptx, as_attachment=True)

    return render_template('index.html', menu_items=menu_items)

if __name__ == '__main__':
    app.run(debug=True)
